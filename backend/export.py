from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, Any
import logging

THEME_COLORS = {
    'minimal-light': {'bg': 'FFFFFF', 'title': '1A1A1A', 'accent': '6366F1', 'text': '374151'},
    'midnight-professional': {'bg': '0F172A', 'title': 'E2E8F0', 'accent': '38BDF8', 'text': 'CBD5E1'},
    'warm-neutral': {'bg': 'FAFAF9', 'title': '292524', 'accent': 'D97706', 'text': '44403C'},
    'forest-academic': {'bg': 'F0F4F0', 'title': '1C2B1A', 'accent': '16A34A', 'text': '1C2B1A'},
    'slate-coral': {'bg': 'F8FAFC', 'title': '1E293B', 'accent': 'F43F5E', 'text': '334155'},
    'monochrome-editorial': {'bg': 'F5F5F5', 'title': '000000', 'accent': '525252', 'text': '262626'},
    'deep-purple-tech': {'bg': '1E1B4B', 'title': 'E9D5FF', 'accent': 'A855F7', 'text': 'C4B5FD'}
}

def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def generate_pptx(deck: Dict[str, Any], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    theme_name = deck.get('theme', 'minimal-light')
    theme = THEME_COLORS.get(theme_name, THEME_COLORS['minimal-light'])
    
    bg_color = _hex_to_rgb(theme['bg'])
    title_color = _hex_to_rgb(theme['title'])
    accent_color = _hex_to_rgb(theme['accent'])
    text_color = _hex_to_rgb(theme['text'])
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in deck.get('slides', []):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        layout_type = slide_data.get('layout', 'content')
        
        if layout_type == 'title':
            _add_title_slide(slide, slide_data, title_color, accent_color, text_color)
        elif layout_type == 'statistics':
            _add_statistics_slide(slide, slide_data, title_color, accent_color, text_color)
        elif layout_type == 'quote':
            _add_quote_slide(slide, slide_data, title_color, accent_color, text_color)
        elif layout_type == 'closing':
            _add_closing_slide(slide, slide_data, accent_color)
        else:
            _add_content_slide(slide, slide_data, title_color, accent_color, text_color)
            
        note_text = slide_data.get('speaker_note')
        if note_text and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = note_text
            
    prs.save(output_path)

def _add_title_slide(slide, data, title_color, accent_color, text_color):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color
    
    bullets = data.get('bullets', [])
    if bullets:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(11), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = bullets[0]
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(6.8), Inches(12), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

def _add_content_slide(slide, data, title_color, accent_color, text_color):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.85))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color
    
    bullets = data.get('bullets', [])
    if bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(11.5), Inches(5.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f'• {b}'
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(10)
            
    stat = data.get('key_stat')
    if stat:
        stat_box = slide.shapes.add_textbox(Inches(10.5), Inches(2), Inches(2.5), Inches(2.5))
        tf = stat_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = str(stat.get('value', ''))
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = str(stat.get('label', ''))
        p2.font.size = Pt(12)
        p2.font.color.rgb = text_color
        p2.alignment = PP_ALIGN.CENTER

def _add_statistics_slide(slide, data, title_color, accent_color, text_color):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.85))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color
    
    stat = data.get('key_stat')
    if stat:
        stat_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.8), Inches(6.5), Inches(2.5))
        tf = stat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(stat.get('value', ''))
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = str(stat.get('label', ''))
        p2.font.size = Pt(20)
        p2.font.color.rgb = text_color
        p2.alignment = PP_ALIGN.CENTER
        
    bullets = data.get('bullets', [])
    if bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12), Inches(2))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets[:3]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f'• {b}'
            p.font.size = Pt(16)
            p.font.color.rgb = text_color

def _add_quote_slide(slide, data, title_color, accent_color, text_color):
    bar = slide.shapes.add_shape(1, Inches(0.3), Inches(1.5), Inches(0.12), Inches(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '\u201c'
    p.font.size = Pt(60)
    p.font.color.rgb = accent_color
    
    quote_text = data.get('bullets', [])[0] if data.get('bullets') else data.get('title', '')
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote_text
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.color.rgb = title_color
    
    if len(data.get('bullets', [])) > 1:
        attr = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
        tf = attr.text_frame
        p = tf.paragraphs[0]
        p.text = '\u2014 ' + data['bullets'][1]
        p.font.size = Pt(16)
        p.font.color.rgb = accent_color

def _add_closing_slide(slide, data, accent_color):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = accent_color
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Thank You')
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    bullets = data.get('bullets', [])
    if bullets:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = bullets[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
